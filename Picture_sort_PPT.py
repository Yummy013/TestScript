# coding: utf-8

# 対象となるフォルダに移動  cd /Users/yt/GoogleDrive/BZ-X710
# →Python Picture_sort_PPT.py ターゲットフォルダ名

# pip install xlsxwriter
# pip install python-pptx  # ppt
# pip install python-docx
# pip install pdfminer.six
# pip install tqdm # progressbarを表示
import pptx
from pptx.util import Inches
from pptx.util import Pt
from pptx import Presentation
import glob
import os # os.getcwd() や os.path.isfile()で使用
from tqdm import tqdm # tqdm progressbarを表示
import time #time.sleepで使用
import sys #sys.argvリストがコマンドライン引数を保持


# Presentationインスタンスの作成
ppt = Presentation()
width = ppt.slide_width # 幅
height = ppt.slide_height # 高さ

#使用するスライドの種類
title_slide_layout = ppt.slide_layouts[0] #Title Slideの作成
bullet_slide_layout = ppt.slide_layouts[1] #Title and Contentの作成
blank_slide_layout = ppt.slide_layouts[6] #Blankの作成

############################################################################
#Title Slide
slide = ppt.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]


import datetime #今日の日付指定
now = datetime.datetime.now() #今日の日付を取得
today = str(now.strftime("%Y/%m/%d")) #str関数でdatetimeオブジェクトを文字列に変換

title.text = sys.argv[1] #フォルダ名をタイトルにするため　.py以降の１つ目の引数にターゲットのフォルダ名を指定
subtitle.text = today

############################################################################
#Title and Contentの作成
slide_description = ppt.slides.add_slide(bullet_slide_layout)
shapes = slide_description.shapes

title_shape = shapes.title
body_shape = shapes.placeholders[1]

#Title and Contentにタイトルとテキストを書き込む
title_shape.text = 'Condition'
tf = body_shape.text_frame
tf.text = 'Drug: Chicago Blue 0.5% in ACSF　0.1μl/min 2min' # sys.argv[2]で２つ目の引数で記入してもよい
tf.paragraphs[0].font.size = Pt(28) #font size
tf.paragraphs[0].font.bold = False #font bold

p = tf.add_paragraph()
p.text = 'Coordinate: BLA (AP-1.30mm　ML±3.30mm)'
tf.paragraphs[1].font.size = Pt(28) # font size
tf.paragraphs[1].font.bold = False # font bold

p = tf.add_paragraph()
p.text = 'Canula size'
tf.paragraphs[1].font.size = Pt(28) # font size
tf.paragraphs[1].font.bold = False # font bold

p = tf.add_paragraph()
p.text = 'Guide canula 26G, pedestal 5mm,Guide length 4.1mm'
p.level = 1  # down the bullet level
p.font.size = Pt(24)  # font size

p = tf.add_paragraph()
p.text = 'Internal canula 33G, pedestal 5mm, Guide length 4.1mm/露出0.5mm'
p.level = 1  # down the bullet levelS
p.font.size = Pt(24)  # font size

############################################################################
# 画像ファイルの読み込み
fnms = sorted(glob.glob(os.getcwd()+ '/' + sys.argv[1] + '/**/*CH*tif', recursive=True))
# 名前でソートしたい場合： sorted(glob.glob('*.png'))
# 引数recursive=Trueとして**を使うと、あらゆるファイルや0個以上のディレクトリおよびサブディレクトリにマッチする。(Python3.5以降)
# 例えば*を使って任意のディレクトリ名にマッチさせると同一階層のファイルしか抽出できないが、**を使うとあらゆる中間ディレクトリに対応してマッチさせることができる。
# カレントディレクトリの絶対パスを取得
#　sys.argvリストがコマンドライン引数を保持

tx_left = tx_top = tx_width = tx_height = Inches(1)
step = 2 #for文で2個ずつ回す

# ファイル毎にループ
for fnm in tqdm(range(0, len (fnms), step)): #tqdm()でprogressbarを表示
    time.sleep(0.01)  #処理を一旦停止しCPUへの負荷を軽減
    fnms[fnm : fnm + step] # https://qiita.com/Akio-1978/items/c003783517df23d360d0

    # 白紙のスライドの追加
    slide = ppt.slides.add_slide(blank_slide_layout)

    # 画像左の挿入
    pic = slide.shapes.add_picture(fnms[fnm], width/4,height/2, width/2, height/2)
    # 右に移動
    pic.left = int(0) #中心にするにはpic.left = int( ( width  - pic.width  ) / 2 )
    pic.top  = int( ( height - pic.height ) / 2 )

    # 画像右の挿入
    pic = slide.shapes.add_picture(fnms[fnm + 1], width/4,height/2, width/2, height/2)
    # 左に移動
    pic.left = int( width  - pic.width )
    pic.top  = int( ( height - pic.height ) / 2 )

    #ついでに空きスペースにテキストボックスを挿入する
    txBox = slide.shapes.add_textbox(tx_left, tx_top, tx_width, tx_height)
    tB = txBox.text_frame
    tB.text = fnms[fnm].strip("/Users/yt/GoogleDrive/BZ-X710/") #rstrip()関数は文字列の右側の文字を除去したコピーを返す     # https://techacademy.jp/magazine/33799
    tB.paragraphs[0].font.size = Pt(24) # font size

############################################################################
# 名前をつけて保存
ppt.save('CannulaLocation_' + sys.argv[1] + '.pptx') #フォルダ名をファイル名にいれる

############################################################################
